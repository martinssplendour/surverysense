from __future__ import annotations

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt

from app.features.export.report_export_service._constants import (
    _PPTX_CONTENT_LEFT,
    _PPTX_CONTENT_WIDTH,
    _PPTX_DETAIL_RGB,
    _PPTX_SLIDE_BACKGROUND_RGB,
    _PPTX_SLIDE_HEIGHT,
    _PPTX_TEXT_RGB,
)
from app.features.export.report_export_service.chart_image import DecodedChartImage


class PptxReportBuilder:
    def __init__(self, *, content_service, chart_image_service) -> None:
        self.content_service = content_service
        self.chart_image_service = chart_image_service

    def build(self, *, request, charts: list[DecodedChartImage]) -> bytes:
        presentation = Presentation()
        self._build_title_slide(presentation, request)

        for chart in charts:
            self._build_chart_slide(presentation, chart)

        self._build_summary_slide(presentation, request)
        self._build_representative_slides(presentation, request)

        output = BytesIO()
        presentation.save(output)
        return output.getvalue()

    def _build_title_slide(self, presentation: Presentation, request) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_slide_background(slide)

        title_top = 3.3
        subtitle_top = 4.02
        title_box_width = _PPTX_CONTENT_WIDTH * 0.76
        title_box_left = _PPTX_CONTENT_LEFT

        title_box = slide.shapes.add_textbox(
            PptxInches(title_box_left),
            PptxInches(title_top),
            PptxInches(title_box_width),
            PptxInches(0.7),
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = self.content_service.build_report_title()
        title_paragraph.font.size = PptxPt(26)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = PptxRGBColor(*_PPTX_TEXT_RGB)
        title_paragraph.alignment = PP_ALIGN.LEFT

        subtitle_box = slide.shapes.add_textbox(
            PptxInches(title_box_left),
            PptxInches(subtitle_top),
            PptxInches(title_box_width),
            PptxInches(0.68),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.word_wrap = True
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.text = self.content_service.build_subtitle(request)
        subtitle_paragraph.font.size = PptxPt(10)
        subtitle_paragraph.font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
        subtitle_paragraph.alignment = PP_ALIGN.LEFT

    def _build_chart_slide(self, presentation: Presentation, chart: DecodedChartImage) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_slide_title(slide, chart.title)

        content_top = 1.0
        if chart.caption:
            caption_box = slide.shapes.add_textbox(
                PptxInches(_PPTX_CONTENT_LEFT),
                PptxInches(1.02),
                PptxInches(_PPTX_CONTENT_WIDTH),
                PptxInches(0.4),
            )
            caption_frame = caption_box.text_frame
            caption_frame.paragraphs[0].text = chart.caption
            caption_frame.paragraphs[0].font.size = PptxPt(10)
            caption_frame.paragraphs[0].font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
            caption_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            content_top = 1.42

        width_inches, height_inches = self.chart_image_service.fit_image_to_bounds(
            width=chart.width,
            height=chart.height,
            max_width=_PPTX_CONTENT_WIDTH,
            max_height=_PPTX_SLIDE_HEIGHT - content_top - 0.45,
        )
        left = max(0.25, _PPTX_CONTENT_LEFT - 0.12)
        top = max(content_top, (_PPTX_SLIDE_HEIGHT - height_inches) / 2)
        slide.shapes.add_picture(
            BytesIO(chart.image_bytes),
            PptxInches(left),
            PptxInches(top),
            width=PptxInches(width_inches),
            height=PptxInches(height_inches),
        )

    def _build_summary_slide(self, presentation: Presentation, request) -> None:
        summary_box_width = _PPTX_CONTENT_WIDTH * 0.7
        summary_box_left = _PPTX_CONTENT_LEFT
        group_sections = self.content_service.build_group_summary_sections(request)
        if group_sections:
            for chunk_start in range(0, min(len(group_sections), 8), 4):
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                self._set_slide_background(slide)
                title = self.content_service.build_summary_heading(request)
                if chunk_start > 0:
                    title = f"{title} (continued)"
                self._add_slide_title(slide, title)

                top = 1.2
                for section in group_sections[chunk_start:chunk_start + 4]:
                    box = slide.shapes.add_textbox(
                        PptxInches(summary_box_left),
                        PptxInches(top),
                        PptxInches(summary_box_width),
                        PptxInches(1.25),
                    )
                    frame = box.text_frame
                    frame.clear()
                    frame.word_wrap = True

                    heading = frame.paragraphs[0]
                    heading.text = section.label
                    heading.font.size = PptxPt(16)
                    heading.font.bold = True
                    heading.font.color.rgb = PptxRGBColor(*_PPTX_TEXT_RGB)
                    heading.alignment = PP_ALIGN.LEFT

                    paragraph = frame.add_paragraph()
                    paragraph.text = section.summary
                    paragraph.font.size = PptxPt(12)
                    paragraph.font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.space_after = PptxPt(6)
                    top += 1.35
            return

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_slide_background(slide)
        self._add_slide_title(slide, self.content_service.build_summary_heading(request))

        box = slide.shapes.add_textbox(
            PptxInches(summary_box_left),
            PptxInches(1.35),
            PptxInches(summary_box_width),
            PptxInches(5.6),
        )
        frame = box.text_frame
        frame.word_wrap = True
        findings = self.content_service.build_summary_lines(request)[:8]
        for index, line in enumerate(findings):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.bullet = True
            paragraph.font.size = PptxPt(14)
            paragraph.font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = PptxPt(8)

    def _build_representative_slides(self, presentation: Presentation, request) -> None:
        sections = self.content_service.build_representative_sections(request)
        if not sections:
            return

        for chunk_start in range(0, len(sections), 2):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            self._set_slide_background(slide)
            self._add_slide_title(slide, self.content_service.build_representative_heading())

            top = 1.25
            group_box_width = _PPTX_CONTENT_WIDTH * 0.7
            group_box_left = _PPTX_CONTENT_LEFT
            for label, examples in sections[chunk_start:chunk_start + 2]:
                group_box = slide.shapes.add_textbox(
                    PptxInches(group_box_left),
                    PptxInches(top),
                    PptxInches(group_box_width),
                    PptxInches(2.65),
                )
                frame = group_box.text_frame
                frame.clear()
                frame.word_wrap = True

                heading = frame.paragraphs[0]
                heading.text = label
                heading.font.size = PptxPt(16)
                heading.font.bold = True
                heading.font.color.rgb = PptxRGBColor(*_PPTX_TEXT_RGB)
                heading.alignment = PP_ALIGN.LEFT

                for index, example in enumerate(examples, start=1):
                    paragraph = frame.add_paragraph()
                    paragraph.text = f"{index}. {example}"
                    paragraph.font.size = PptxPt(12)
                    paragraph.font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.space_after = PptxPt(6)
                top += 2.8

    @staticmethod
    def _set_slide_background(slide) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = PptxRGBColor(*_PPTX_SLIDE_BACKGROUND_RGB)

    @staticmethod
    def _add_slide_title(slide, title: str) -> None:
        title_box = slide.shapes.add_textbox(
            PptxInches(_PPTX_CONTENT_LEFT),
            PptxInches(0.45),
            PptxInches(_PPTX_CONTENT_WIDTH),
            PptxInches(0.55),
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        paragraph = title_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = PptxPt(20)
        paragraph.font.bold = True
        paragraph.font.color.rgb = PptxRGBColor(*_PPTX_TEXT_RGB)
        paragraph.alignment = PP_ALIGN.LEFT
