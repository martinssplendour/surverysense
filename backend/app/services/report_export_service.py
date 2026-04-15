from __future__ import annotations

import base64
import re
from dataclasses import dataclass
from io import BytesIO

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image as ReportLabImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)

from app.models.api import AnalysisExportChartModel, AnalysisExportRequest

_DATA_URL_PATTERN = re.compile(r"^data:image/(?P<kind>[a-zA-Z0-9.+-]+);base64,(?P<data>.+)$")
_FILENAME_PATTERN = re.compile(r"[^a-z0-9]+")
_REPORT_TITLE_COLOR = "#2a3f5f"
_REPORT_TITLE_RGB = (42, 63, 95)


@dataclass(slots=True)
class ExportedReportArtifact:
    filename: str
    media_type: str
    content: bytes


@dataclass(slots=True)
class DecodedChartImage:
    title: str
    caption: str
    image_bytes: bytes
    width: int
    height: int


class AnalysisReportExportService:
    def export_report(self, *, result_id: str, request: AnalysisExportRequest) -> ExportedReportArtifact:
        if request.analysis_result.result_id != result_id:
            raise ValueError("The report payload does not match the requested analysis result.")

        charts = [self._decode_chart(chart) for chart in request.charts]
        content: bytes
        media_type: str
        extension = request.format

        if request.format == "pdf":
            content = self._build_pdf(request=request, charts=charts)
            media_type = "application/pdf"
        elif request.format == "docx":
            content = self._build_docx(request=request, charts=charts)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        else:
            content = self._build_pptx(request=request, charts=charts)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        stem = self._slugify(self._build_filename_stem(request))
        return ExportedReportArtifact(
            filename=f"{stem}.{extension}",
            media_type=media_type,
            content=content,
        )

    def _build_pdf(self, *, request: AnalysisExportRequest, charts: list[DecodedChartImage]) -> bytes:
        buffer = BytesIO()
        document = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            leftMargin=42,
            rightMargin=42,
            topMargin=42,
            bottomMargin=42,
        )
        styles = self._build_pdf_styles()
        story: list[object] = []

        story.append(Paragraph(self._escape(self._build_report_title()), styles["title"]))
        story.append(Paragraph(self._escape(self._build_subtitle(request)), styles["subtitle"]))
        story.append(Spacer(1, 18))

        filters_text = self._filters_text(request)
        if filters_text:
            story.append(Paragraph(self._escape(filters_text), styles["body"]))
            story.append(Spacer(1, 12))

        if charts:
            for index, chart in enumerate(charts):
                story.append(Paragraph(self._escape(chart.title), styles["chart_title"]))
                if chart.caption:
                    story.append(Paragraph(self._escape(chart.caption), styles["muted"]))
                story.append(Spacer(1, 8))
                story.append(self._build_pdf_chart_image(chart, max_width=document.width, max_height=320))
                if index < len(charts) - 1:
                    story.append(Spacer(1, 18))
            story.append(PageBreak())

        story.append(Paragraph(self._build_summary_heading(request), styles["section"]))
        story.append(Spacer(1, 8))
        findings = self._build_summary_lines(request)
        for line in findings:
            story.append(Paragraph(f"• {self._escape(line)}", styles["body"]))
            story.append(Spacer(1, 5))

        representative_sections = self._build_representative_sections(request)
        if representative_sections:
            story.append(Spacer(1, 10))
            story.append(Paragraph("Representative documents", styles["section"]))
            story.append(Spacer(1, 8))
            for label, examples in representative_sections:
                story.append(Paragraph(self._escape(label), styles["chart_title"]))
                for index, example in enumerate(examples, start=1):
                    story.append(Paragraph(f"{index}. {self._escape(example)}", styles["body"]))
                    story.append(Spacer(1, 4))
                story.append(Spacer(1, 8))

        document.build(story, onFirstPage=self._decorate_pdf_page, onLaterPages=self._decorate_pdf_page)
        return buffer.getvalue()

    def _build_docx(self, *, request: AnalysisExportRequest, charts: list[DecodedChartImage]) -> bytes:
        document = Document()
        section = document.sections[0]
        section.top_margin = DocxInches(0.55)
        section.bottom_margin = DocxInches(0.55)
        section.left_margin = DocxInches(0.65)
        section.right_margin = DocxInches(0.65)

        title = document.add_paragraph()
        title.style = document.styles["Title"]
        title.alignment = WD_ALIGN_PARAGRAPH.LEFT
        title_run = title.add_run(self._build_report_title())
        title_run.font.name = "Aptos Display"
        title_run.font.size = DocxPt(22)
        title_run.font.color.rgb = RGBColor(*_REPORT_TITLE_RGB)

        subtitle = document.add_paragraph(self._build_subtitle(request))
        subtitle.style = document.styles["Subtitle"]
        for run in subtitle.runs:
            run.font.name = "Aptos"
            run.font.size = DocxPt(10)
        document.add_paragraph()

        filters_text = self._filters_text(request)
        if filters_text:
            paragraph = document.add_paragraph()
            filters_run = paragraph.add_run(filters_text)
            filters_run.font.name = "Aptos"
            filters_run.font.size = DocxPt(10)

        if charts:
            for chart in charts:
                document.add_heading(chart.title, level=2)
                if chart.caption:
                    caption = document.add_paragraph(chart.caption)
                    for run in caption.runs:
                        run.italic = True
                        run.font.name = "Aptos"
                        run.font.size = DocxPt(9)
                        run.font.color.rgb = RGBColor(98, 91, 82)
                document.add_picture(BytesIO(chart.image_bytes), width=DocxInches(6.45))

        document.add_heading(self._build_summary_heading(request), level=1)
        for line in self._build_summary_lines(request):
            paragraph = document.add_paragraph()
            run = paragraph.add_run(line)
            run.font.name = "Aptos"
            run.font.size = DocxPt(10)

        representative_sections = self._build_representative_sections(request)
        if representative_sections:
            document.add_heading("Representative documents", level=1)
            for label, examples in representative_sections:
                group_heading = document.add_paragraph()
                group_run = group_heading.add_run(label)
                group_run.bold = True
                group_run.font.name = "Aptos"
                group_run.font.size = DocxPt(11)
                for index, example in enumerate(examples, start=1):
                    paragraph = document.add_paragraph()
                    run = paragraph.add_run(f"{index}. {example}")
                    run.font.name = "Aptos"
                    run.font.size = DocxPt(10)

        output = BytesIO()
        document.save(output)
        return output.getvalue()

    def _build_pptx(self, *, request: AnalysisExportRequest, charts: list[DecodedChartImage]) -> bytes:
        presentation = Presentation()
        self._build_pptx_title_slide(presentation, request)

        for chart in charts:
            self._build_pptx_chart_slide(presentation, chart)

        self._build_pptx_summary_slide(presentation, request)
        self._build_pptx_representative_slides(presentation, request)

        output = BytesIO()
        presentation.save(output)
        return output.getvalue()

    def _build_pptx_title_slide(self, presentation: Presentation, request: AnalysisExportRequest) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = PptxRGBColor(247, 242, 234)

        title_box = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(0.8), PptxInches(11.4), PptxInches(1.2))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = self._build_report_title()
        title_paragraph.font.size = PptxPt(24)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = PptxRGBColor(*_REPORT_TITLE_RGB)

        subtitle_box = slide.shapes.add_textbox(PptxInches(0.72), PptxInches(1.95), PptxInches(11.2), PptxInches(1.0))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.word_wrap = True
        subtitle_frame.paragraphs[0].text = self._build_subtitle(request)
        subtitle_frame.paragraphs[0].font.size = PptxPt(12)
        subtitle_frame.paragraphs[0].font.color.rgb = PptxRGBColor(84, 77, 69)
        filters_text = self._filters_text(request)
        if filters_text:
            filters_box = slide.shapes.add_textbox(PptxInches(0.72), PptxInches(2.5), PptxInches(11.0), PptxInches(0.8))
            filters_frame = filters_box.text_frame
            filters_frame.word_wrap = True
            filters_frame.paragraphs[0].text = filters_text
            filters_frame.paragraphs[0].font.size = PptxPt(10)
            filters_frame.paragraphs[0].font.color.rgb = PptxRGBColor(84, 77, 69)

    def _build_pptx_chart_slide(self, presentation: Presentation, chart: DecodedChartImage) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_pptx_slide_background(slide)
        self._add_pptx_slide_title(slide, chart.title)

        if chart.caption:
            caption_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(1.25), PptxInches(11.0), PptxInches(0.45))
            caption_frame = caption_box.text_frame
            caption_frame.paragraphs[0].text = chart.caption
            caption_frame.paragraphs[0].font.size = PptxPt(10)
            caption_frame.paragraphs[0].font.color.rgb = PptxRGBColor(98, 91, 82)

        image = Image.open(BytesIO(chart.image_bytes))
        width_inches, height_inches = self._fit_image_to_bounds(
            width=image.width,
            height=image.height,
            max_width=11.0,
            max_height=5.3,
        )
        left = (13.333 - width_inches) / 2
        top = 1.65 + ((5.3 - height_inches) / 2)
        slide.shapes.add_picture(
            BytesIO(chart.image_bytes),
            PptxInches(left),
            PptxInches(top),
            width=PptxInches(width_inches),
            height=PptxInches(height_inches),
        )

    def _build_pptx_summary_slide(self, presentation: Presentation, request: AnalysisExportRequest) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_pptx_slide_background(slide)
        self._add_pptx_slide_title(slide, self._build_summary_heading(request))

        box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(1.35), PptxInches(11.3), PptxInches(5.6))
        frame = box.text_frame
        frame.word_wrap = True
        findings = self._build_summary_lines(request)[:8]
        for index, line in enumerate(findings):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.bullet = True
            paragraph.font.size = PptxPt(14)
            paragraph.font.color.rgb = PptxRGBColor(61, 53, 45)
            paragraph.space_after = PptxPt(8)

    def _build_pptx_representative_slides(self, presentation: Presentation, request: AnalysisExportRequest) -> None:
        sections = self._build_representative_sections(request)
        if not sections:
            return

        for chunk_start in range(0, len(sections), 2):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            self._set_pptx_slide_background(slide)
            self._add_pptx_slide_title(slide, "Representative documents")

            top = 1.25
            for label, examples in sections[chunk_start:chunk_start + 2]:
                group_box = slide.shapes.add_textbox(PptxInches(0.8), PptxInches(top), PptxInches(11.3), PptxInches(2.6))
                frame = group_box.text_frame
                frame.word_wrap = True
                heading = frame.paragraphs[0]
                heading.text = label
                heading.font.size = PptxPt(16)
                heading.font.bold = True
                heading.font.color.rgb = PptxRGBColor(41, 75, 59)
                for index, example in enumerate(examples, start=1):
                    paragraph = frame.add_paragraph()
                    paragraph.text = f"{index}. {example}"
                    paragraph.font.size = PptxPt(12)
                    paragraph.font.color.rgb = PptxRGBColor(61, 53, 45)
                    paragraph.space_after = PptxPt(6)
                top += 2.8

    def _build_pdf_styles(self) -> dict[str, ParagraphStyle]:
        base = getSampleStyleSheet()
        return {
            "title": ParagraphStyle(
                "ReportTitle",
                parent=base["Title"],
                fontName="Helvetica-Bold",
                fontSize=22,
                leading=26,
                textColor=colors.HexColor(_REPORT_TITLE_COLOR),
                alignment=TA_LEFT,
                spaceAfter=6,
            ),
            "subtitle": ParagraphStyle(
                "ReportSubtitle",
                parent=base["BodyText"],
                fontName="Helvetica",
                fontSize=10,
                leading=14,
                textColor=colors.HexColor("#5e574f"),
                spaceAfter=0,
            ),
            "section": ParagraphStyle(
                "ReportSection",
                parent=base["Heading2"],
                fontName="Helvetica-Bold",
                fontSize=14,
                textColor=colors.HexColor("#294b3b"),
                spaceAfter=4,
            ),
            "chart_title": ParagraphStyle(
                "ReportChartTitle",
                parent=base["Heading3"],
                fontName="Helvetica-Bold",
                fontSize=11,
                textColor=colors.HexColor("#3d352d"),
                spaceAfter=2,
            ),
            "muted": ParagraphStyle(
                "ReportMuted",
                parent=base["BodyText"],
                fontName="Helvetica-Oblique",
                fontSize=9,
                leading=12,
                textColor=colors.HexColor("#6d655b"),
            ),
            "body": ParagraphStyle(
                "ReportBody",
                parent=base["BodyText"],
                fontName="Helvetica",
                fontSize=10,
                leading=14,
                textColor=colors.HexColor("#3d352d"),
            ),
        }

    def _build_pdf_chart_image(self, chart: DecodedChartImage, *, max_width: float, max_height: float) -> ReportLabImage:
        width_inches, height_inches = self._fit_image_to_bounds(
            width=chart.width,
            height=chart.height,
            max_width=max_width / inch,
            max_height=max_height / inch,
        )
        return ReportLabImage(BytesIO(chart.image_bytes), width=width_inches * inch, height=height_inches * inch)

    def _decorate_pdf_page(self, canvas, document) -> None:
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor("#d8cdbf"))
        canvas.setFillColor(colors.HexColor("#294b3b"))
        canvas.setFont("Helvetica", 9)
        canvas.drawRightString(document.pagesize[0] - 42, 20, f"Page {document.page}")
        canvas.restoreState()

    def _build_summary_lines(self, request: AnalysisExportRequest) -> list[str]:
        result = request.analysis_result
        if result.ngram_buckets:
            findings: list[str] = []
            for bucket in result.ngram_buckets:
                top_items = bucket.items[:5]
                if not top_items:
                    continue
                lines = ", ".join(
                    f"{item.term} ({item.document_count} responses)"
                    for item in top_items
                )
                findings.append(f"{bucket.label}: {lines}")
            return findings or ["No phrase-level findings were available for export."]

        if result.groups:
            ordered = sorted(result.groups, key=lambda group: (-group.count, group.label))
            findings = []
            for group in ordered[:8]:
                share_pct = round(group.share * 100)
                terms = ", ".join(group.terms[:4]) if group.terms else "no top terms available"
                findings.append(
                    f"{group.label}: {group.count} responses ({share_pct}%). Top terms: {terms}."
                )
            return findings or ["No grouped findings were available for export."]

        return ["The selected analysis completed without exportable grouped findings."]

    def _build_subtitle(self, request: AnalysisExportRequest) -> str:
        return self._display_column_label(request.analysis_result.text_column_name)

    def _build_report_title(self) -> str:
        return "Verbatim Analysis Report"

    def _build_summary_heading(self, request: AnalysisExportRequest) -> str:
        if request.analysis_result.ngram_buckets:
            return "Phrase summaries"
        return "Group summaries"

    def _build_representative_sections(self, request: AnalysisExportRequest) -> list[tuple[str, list[str]]]:
        if not request.analysis_result.groups:
            return []

        sections: list[tuple[str, list[str]]] = []
        ordered_groups = sorted(request.analysis_result.groups, key=lambda group: (-group.count, group.label))
        for group in ordered_groups:
            examples = [
                self._truncate_text(example.text, limit=240)
                for example in group.examples[:3]
                if example.text.strip()
            ]
            if not examples:
                continue
            sections.append((group.label, examples))
        return sections

    def _build_filename_stem(self, request: AnalysisExportRequest) -> str:
        source_stem = self._strip_extension(request.source_filename or "analysis")
        method = request.analysis_result.model_label.lower().replace(" ", "-")
        return f"{source_stem}-{method}-report"

    def _filters_text(self, request: AnalysisExportRequest) -> str:
        if not request.active_filters:
            return ""
        return " | ".join(
            f"{item.display_name or item.column_name}: {', '.join(item.values)}"
            for item in request.active_filters
            if item.values
        )

    def _decode_chart(self, chart: AnalysisExportChartModel) -> DecodedChartImage:
        match = _DATA_URL_PATTERN.match(chart.image_data_url.strip())
        if match is None:
            raise ValueError(f"Chart '{chart.title}' does not contain a supported image data URL.")

        image_bytes = base64.b64decode(match.group("data"))
        image = Image.open(BytesIO(image_bytes))
        image.load()
        normalized = BytesIO()
        if image.mode not in {"RGB", "RGBA"}:
            image = image.convert("RGB")
        image.save(normalized, format="PNG")
        return DecodedChartImage(
            title=chart.title.strip() or "Chart",
            caption=self._sanitize_chart_caption(chart.caption),
            image_bytes=normalized.getvalue(),
            width=int(image.width),
            height=int(image.height),
        )

    def _fit_image_to_bounds(self, *, width: int, height: int, max_width: float, max_height: float) -> tuple[float, float]:
        if width <= 0 or height <= 0:
            return max_width, min(max_height, max_width * 0.6)
        scale = min(max_width / float(width), max_height / float(height))
        return width * scale, height * scale

    def _display_column_label(self, value: str) -> str:
        return re.sub(r"__idx_\d+$", "", value or "")

    def _strip_extension(self, filename: str) -> str:
        value = (filename or "").strip()
        if "." not in value:
            return value
        return value.rsplit(".", 1)[0]

    def _slugify(self, value: str) -> str:
        lowered = value.strip().lower()
        normalized = _FILENAME_PATTERN.sub("-", lowered).strip("-")
        return normalized or "analysis-report"

    def _escape(self, value: str) -> str:
        return (
            str(value)
            .replace("&", "&amp;")
            .replace("<", "&lt;")
            .replace(">", "&gt;")
        )

    def _truncate_text(self, value: str, *, limit: int) -> str:
        normalized = " ".join((value or "").split())
        if len(normalized) <= limit:
            return normalized
        return normalized[: max(0, limit - 3)].rstrip() + "..."

    def _sanitize_chart_caption(self, value: str | None) -> str:
        caption = " ".join((value or "").split())
        lower_caption = caption.casefold()
        interactive_phrases = (
            "hover to see",
            "click a bar",
            "matching responses",
            "matching groups responses",
            "matching themes responses",
        )
        if any(phrase in lower_caption for phrase in interactive_phrases):
            return ""
        return caption

    def _set_pptx_slide_background(self, slide) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = PptxRGBColor(247, 242, 234)

    def _add_pptx_slide_title(self, slide, title: str) -> None:
        title_box = slide.shapes.add_textbox(PptxInches(0.7), PptxInches(0.45), PptxInches(11.4), PptxInches(0.55))
        title_frame = title_box.text_frame
        title_frame.clear()
        paragraph = title_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = PptxPt(22)
        paragraph.font.bold = True
        paragraph.font.color.rgb = PptxRGBColor(41, 75, 59)
