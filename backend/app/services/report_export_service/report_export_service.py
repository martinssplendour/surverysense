from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches as DocxInches
from docx.shared import Pt as DocxPt
from docx.shared import RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptxRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptxInches
from pptx.util import Pt as PptxPt
from reportlab.lib import colors
from reportlab.lib.enums import TA_LEFT
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image as ReportLabImage,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
)

from app.models.api import AnalysisExportChartModel, AnalysisExportRequest
from app.services.report_export_service._constants import (
    _PPTX_CONTENT_LEFT,
    _PPTX_CONTENT_WIDTH,
    _PPTX_DETAIL_RGB,
    _PPTX_SLIDE_BACKGROUND_RGB,
    _PPTX_SLIDE_HEIGHT,
    _PPTX_TEXT_RGB,
    _REPORT_TITLE_COLOR,
    _REPORT_TITLE_RGB,
)
from app.services.report_export_service.chart_image_service import (
    DecodedChartImage,
    ReportChartImageService,
)
from app.services.report_export_service.report_content_service import (
    GroupSummarySection,
    ReportContentService,
)


@dataclass(slots=True)
class ExportedReportArtifact:
    filename: str
    media_type: str
    content: bytes


class AnalysisReportExportService:
    def __init__(self, result_store_service=None) -> None:
        self.result_store_service = result_store_service
        self.content_service = ReportContentService(result_store_service=result_store_service)
        self.chart_image_service = ReportChartImageService(
            sanitize_chart_caption=self.content_service.sanitize_chart_caption,
        )

    def export_report(self, *, result_id: str, request: AnalysisExportRequest) -> ExportedReportArtifact:
        if request.analysis_result.result_id != result_id:
            raise ValueError("The report payload does not match the requested analysis result.")

        charts = [self._normalize_export_chart_image(self._decode_chart(chart)) for chart in request.charts]
        if request.format == "pdf":
            content = self._build_pdf(request=request, charts=charts)
            media_type = "application/pdf"
        elif request.format == "docx":
            content = self._build_docx(request=request, charts=charts)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        else:
            content = self._build_pptx(request=request, charts=charts)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

        stem = self._slugify(self._build_filename_stem(request))
        return ExportedReportArtifact(
            filename=f"{stem}.{request.format}",
            media_type=media_type,
            content=content,
        )

    def _build_pdf(self, *, request: AnalysisExportRequest, charts: list[DecodedChartImage]) -> bytes:
        buffer = BytesIO()
        document = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            leftMargin=42,
            rightMargin=42,
            topMargin=42,
            bottomMargin=42,
        )
        styles = self._build_pdf_styles()
        story: list[object] = []

        story.append(Paragraph(self._escape(self._build_report_title()), styles["title"]))
        story.append(Paragraph(self._escape(self._build_subtitle(request)), styles["subtitle"]))
        story.append(Spacer(1, 12))

        if charts:
            for index, chart in enumerate(charts):
                story.append(Paragraph(self._escape(chart.title), styles["plot_title"]))
                if chart.caption:
                    story.append(Paragraph(self._escape(chart.caption), styles["plot_caption"]))
                story.append(Spacer(1, 8))
                story.append(self._build_pdf_chart_image(chart, max_width=document.width, max_height=320))
                if index < len(charts) - 1:
                    story.append(Spacer(1, 18))
            story.append(PageBreak())

        story.append(Paragraph(self._build_summary_heading(request), styles["section"]))
        story.append(Spacer(1, 8))
        group_sections = self._build_group_summary_sections(request)
        if group_sections:
            for section in group_sections[:8]:
                story.append(Paragraph(self._escape(section.label), styles["chart_title"]))
                story.append(Paragraph(self._escape(section.summary), styles["body"]))
                story.append(Spacer(1, 6))
        else:
            for line in self._build_summary_lines(request):
                story.append(Paragraph(f"&#8226; {self._escape(line)}", styles["body"]))
                story.append(Spacer(1, 5))

        representative_sections = self._build_representative_sections(request)
        if representative_sections:
            story.append(Spacer(1, 10))
            story.append(Paragraph(self._build_representative_heading(), styles["section"]))
            story.append(Spacer(1, 8))
            for label, examples in representative_sections:
                story.append(Paragraph(self._escape(label), styles["chart_title"]))
                for index, example in enumerate(examples, start=1):
                    story.append(Paragraph(f"{index}. {self._escape(example)}", styles["body"]))
                    story.append(Spacer(1, 4))
                story.append(Spacer(1, 8))

        document.build(story, onFirstPage=self._decorate_pdf_page, onLaterPages=self._decorate_pdf_page)
        return buffer.getvalue()

    def _build_docx(self, *, request: AnalysisExportRequest, charts: list[DecodedChartImage]) -> bytes:
        document = Document()
        section = document.sections[0]
        section.top_margin = DocxInches(0.55)
        section.bottom_margin = DocxInches(0.55)
        section.left_margin = DocxInches(0.65)
        section.right_margin = DocxInches(0.65)

        title = document.add_paragraph()
        title.style = document.styles["Title"]
        title.alignment = WD_ALIGN_PARAGRAPH.LEFT
        title_run = title.add_run(self._build_report_title())
        title_run.font.name = "Aptos Display"
        title_run.font.size = DocxPt(22)
        title_run.font.color.rgb = RGBColor(*_REPORT_TITLE_RGB)

        metadata = document.add_paragraph(self._build_subtitle(request))
        metadata.style = document.styles["Subtitle"]
        for run in metadata.runs:
            run.font.name = "Aptos"
            run.font.size = DocxPt(10)

        document.add_paragraph()

        if charts:
            for chart in charts:
                heading = document.add_paragraph()
                heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
                heading_run = heading.add_run(chart.title)
                heading_run.bold = True
                heading_run.font.name = "Aptos Display"
                heading_run.font.size = DocxPt(14)
                heading_run.font.color.rgb = RGBColor(*_REPORT_TITLE_RGB)
                if chart.caption:
                    caption = document.add_paragraph(chart.caption)
                    caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
                    for run in caption.runs:
                        run.italic = True
                        run.font.name = "Aptos"
                        run.font.size = DocxPt(9)
                        run.font.color.rgb = RGBColor(98, 91, 82)
                image_paragraph = document.add_paragraph()
                image_paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                image_paragraph.add_run().add_picture(BytesIO(chart.image_bytes), width=DocxInches(6.45))

        document.add_heading(self._build_summary_heading(request), level=1)
        group_sections = self._build_group_summary_sections(request)
        if group_sections:
            for section in group_sections[:8]:
                heading = document.add_paragraph()
                heading_run = heading.add_run(section.label)
                heading_run.bold = True
                heading_run.font.name = "Aptos"
                heading_run.font.size = DocxPt(11)

                paragraph = document.add_paragraph()
                run = paragraph.add_run(section.summary)
                run.font.name = "Aptos"
                run.font.size = DocxPt(10)
        else:
            for line in self._build_summary_lines(request):
                paragraph = document.add_paragraph()
                run = paragraph.add_run(line)
                run.font.name = "Aptos"
                run.font.size = DocxPt(10)

        representative_sections = self._build_representative_sections(request)
        if representative_sections:
            document.add_heading(self._build_representative_heading(), level=1)
            for label, examples in representative_sections:
                group_heading = document.add_paragraph()
                group_run = group_heading.add_run(label)
                group_run.bold = True
                group_run.font.name = "Aptos"
                group_run.font.size = DocxPt(11)
                for index, example in enumerate(examples, start=1):
                    paragraph = document.add_paragraph()
                    run = paragraph.add_run(f"{index}. {example}")
                    run.font.name = "Aptos"
                    run.font.size = DocxPt(10)

        output = BytesIO()
        document.save(output)
        return output.getvalue()

    def _build_pptx(self, *, request: AnalysisExportRequest, charts: list[DecodedChartImage]) -> bytes:
        presentation = Presentation()
        self._build_pptx_title_slide(presentation, request)

        for chart in charts:
            self._build_pptx_chart_slide(presentation, chart)

        self._build_pptx_summary_slide(presentation, request)
        self._build_pptx_representative_slides(presentation, request)

        output = BytesIO()
        presentation.save(output)
        return output.getvalue()

    def _build_pptx_title_slide(self, presentation: Presentation, request: AnalysisExportRequest) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_pptx_slide_background(slide)

        title_top = 3.3
        subtitle_top = 4.02
        title_box_width = _PPTX_CONTENT_WIDTH * 0.76
        title_box_left = _PPTX_CONTENT_LEFT

        title_box = slide.shapes.add_textbox(
            PptxInches(title_box_left),
            PptxInches(title_top),
            PptxInches(title_box_width),
            PptxInches(0.7),
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_paragraph = title_frame.paragraphs[0]
        title_paragraph.text = self._build_report_title()
        title_paragraph.font.size = PptxPt(26)
        title_paragraph.font.bold = True
        title_paragraph.font.color.rgb = PptxRGBColor(*_PPTX_TEXT_RGB)
        title_paragraph.alignment = PP_ALIGN.LEFT

        subtitle_box = slide.shapes.add_textbox(
            PptxInches(title_box_left),
            PptxInches(subtitle_top),
            PptxInches(title_box_width),
            PptxInches(0.68),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.word_wrap = True
        subtitle_paragraph = subtitle_frame.paragraphs[0]
        subtitle_paragraph.text = self._build_subtitle(request)
        subtitle_paragraph.font.size = PptxPt(10)
        subtitle_paragraph.font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
        subtitle_paragraph.alignment = PP_ALIGN.LEFT

    def _build_pptx_chart_slide(self, presentation: Presentation, chart: DecodedChartImage) -> None:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_pptx_slide_background(slide)
        self._add_pptx_slide_title(slide, chart.title)

        content_top = 1.0
        if chart.caption:
            caption_box = slide.shapes.add_textbox(
                PptxInches(_PPTX_CONTENT_LEFT),
                PptxInches(1.02),
                PptxInches(_PPTX_CONTENT_WIDTH),
                PptxInches(0.4),
            )
            caption_frame = caption_box.text_frame
            caption_frame.paragraphs[0].text = chart.caption
            caption_frame.paragraphs[0].font.size = PptxPt(10)
            caption_frame.paragraphs[0].font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
            caption_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
            content_top = 1.42

        width_inches, height_inches = self._fit_image_to_bounds(
            width=chart.width,
            height=chart.height,
            max_width=_PPTX_CONTENT_WIDTH * 0.8,
            max_height=_PPTX_SLIDE_HEIGHT - content_top - 0.45,
        )
        width_inches *= 0.9
        height_inches *= 0.9
        left = max(0.25, _PPTX_CONTENT_LEFT - 0.12)
        top = max(content_top, (_PPTX_SLIDE_HEIGHT - height_inches) / 2)
        slide.shapes.add_picture(
            BytesIO(chart.image_bytes),
            PptxInches(left),
            PptxInches(top),
            width=PptxInches(width_inches),
            height=PptxInches(height_inches),
        )

    def _build_pptx_summary_slide(self, presentation: Presentation, request: AnalysisExportRequest) -> None:
        summary_box_width = _PPTX_CONTENT_WIDTH * 0.7
        summary_box_left = _PPTX_CONTENT_LEFT
        group_sections = self._build_group_summary_sections(request)
        if group_sections:
            for chunk_start in range(0, min(len(group_sections), 8), 4):
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                self._set_pptx_slide_background(slide)
                title = self._build_summary_heading(request)
                if chunk_start > 0:
                    title = f"{title} (continued)"
                self._add_pptx_slide_title(slide, title)

                top = 1.2
                for section in group_sections[chunk_start:chunk_start + 4]:
                    box = slide.shapes.add_textbox(
                        PptxInches(summary_box_left),
                        PptxInches(top),
                        PptxInches(summary_box_width),
                        PptxInches(1.25),
                    )
                    frame = box.text_frame
                    frame.clear()
                    frame.word_wrap = True

                    heading = frame.paragraphs[0]
                    heading.text = section.label
                    heading.font.size = PptxPt(16)
                    heading.font.bold = True
                    heading.font.color.rgb = PptxRGBColor(*_PPTX_TEXT_RGB)
                    heading.alignment = PP_ALIGN.LEFT

                    paragraph = frame.add_paragraph()
                    paragraph.text = section.summary
                    paragraph.font.size = PptxPt(12)
                    paragraph.font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.space_after = PptxPt(6)
                    top += 1.35
            return

        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        self._set_pptx_slide_background(slide)
        self._add_pptx_slide_title(slide, self._build_summary_heading(request))

        box = slide.shapes.add_textbox(
            PptxInches(summary_box_left),
            PptxInches(1.35),
            PptxInches(summary_box_width),
            PptxInches(5.6),
        )
        frame = box.text_frame
        frame.word_wrap = True
        findings = self._build_summary_lines(request)[:8]
        for index, line in enumerate(findings):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.bullet = True
            paragraph.font.size = PptxPt(14)
            paragraph.font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
            paragraph.alignment = PP_ALIGN.LEFT
            paragraph.space_after = PptxPt(8)

    def _build_pptx_representative_slides(self, presentation: Presentation, request: AnalysisExportRequest) -> None:
        sections = self._build_representative_sections(request)
        if not sections:
            return

        for chunk_start in range(0, len(sections), 2):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            self._set_pptx_slide_background(slide)
            self._add_pptx_slide_title(slide, self._build_representative_heading())

            top = 1.25
            group_box_width = _PPTX_CONTENT_WIDTH * 0.7
            group_box_left = _PPTX_CONTENT_LEFT
            for label, examples in sections[chunk_start:chunk_start + 2]:
                group_box = slide.shapes.add_textbox(
                    PptxInches(group_box_left),
                    PptxInches(top),
                    PptxInches(group_box_width),
                    PptxInches(2.65),
                )
                frame = group_box.text_frame
                frame.clear()
                frame.word_wrap = True

                heading = frame.paragraphs[0]
                heading.text = label
                heading.font.size = PptxPt(16)
                heading.font.bold = True
                heading.font.color.rgb = PptxRGBColor(*_PPTX_TEXT_RGB)
                heading.alignment = PP_ALIGN.LEFT

                for index, example in enumerate(examples, start=1):
                    paragraph = frame.add_paragraph()
                    paragraph.text = f"{index}. {example}"
                    paragraph.font.size = PptxPt(12)
                    paragraph.font.color.rgb = PptxRGBColor(*_PPTX_DETAIL_RGB)
                    paragraph.alignment = PP_ALIGN.LEFT
                    paragraph.space_after = PptxPt(6)
                top += 2.8

    def _build_pdf_styles(self) -> dict[str, ParagraphStyle]:
        base = getSampleStyleSheet()
        return {
            "title": ParagraphStyle(
                "ReportTitle",
                parent=base["Title"],
                fontName="Helvetica-Bold",
                fontSize=22,
                leading=26,
                textColor=colors.HexColor(_REPORT_TITLE_COLOR),
                alignment=TA_LEFT,
                spaceAfter=6,
            ),
            "subtitle": ParagraphStyle(
                "ReportSubtitle",
                parent=base["BodyText"],
                fontName="Helvetica",
                fontSize=10,
                leading=14,
                textColor=colors.HexColor("#5e574f"),
                spaceAfter=0,
            ),
            "section": ParagraphStyle(
                "ReportSection",
                parent=base["Heading2"],
                fontName="Helvetica-Bold",
                fontSize=14,
                textColor=colors.HexColor("#294b3b"),
                spaceAfter=4,
            ),
            "chart_title": ParagraphStyle(
                "ReportChartTitle",
                parent=base["Heading3"],
                fontName="Helvetica-Bold",
                fontSize=11,
                textColor=colors.HexColor("#3d352d"),
                spaceAfter=2,
            ),
            "plot_title": ParagraphStyle(
                "ReportPlotTitle",
                parent=base["Heading3"],
                fontName="Helvetica-Bold",
                fontSize=11,
                textColor=colors.HexColor("#3d352d"),
                spaceAfter=2,
            ),
            "plot_caption": ParagraphStyle(
                "ReportPlotCaption",
                parent=base["BodyText"],
                fontName="Helvetica-Oblique",
                fontSize=9,
                leading=12,
                textColor=colors.HexColor("#6d655b"),
                spaceAfter=6,
            ),
            "body": ParagraphStyle(
                "ReportBody",
                parent=base["BodyText"],
                fontName="Helvetica",
                fontSize=10,
                leading=14,
                textColor=colors.HexColor("#3d352d"),
            ),
        }

    def _build_pdf_chart_image(self, chart: DecodedChartImage, *, max_width: float, max_height: float) -> ReportLabImage:
        width_inches, height_inches = self._fit_image_to_bounds(
            width=chart.width,
            height=chart.height,
            max_width=max_width / inch,
            max_height=max_height / inch,
        )
        return ReportLabImage(BytesIO(chart.image_bytes), width=width_inches * inch, height=height_inches * inch)

    def _decorate_pdf_page(self, canvas, document) -> None:
        canvas.saveState()
        canvas.setStrokeColor(colors.HexColor("#d8cdbf"))
        canvas.setFillColor(colors.HexColor("#294b3b"))
        canvas.setFont("Helvetica", 9)
        canvas.drawRightString(document.pagesize[0] - 42, 20, f"Page {document.page}")
        canvas.restoreState()

    def _build_summary_lines(self, request: AnalysisExportRequest) -> list[str]:
        return self.content_service.build_summary_lines(request)

    def _build_group_summary_sections(self, request: AnalysisExportRequest) -> list[GroupSummarySection]:
        return self.content_service.build_group_summary_sections(request)

    def _get_export_groups(self, request: AnalysisExportRequest):
        return self.content_service.get_export_groups(request)

    def _build_active_filter_lookup(self, request: AnalysisExportRequest) -> dict[str, list[str]]:
        return self.content_service.build_active_filter_lookup(request)

    def _build_subtitle(self, request: AnalysisExportRequest) -> str:
        return self.content_service.build_subtitle(request)

    def _build_report_title(self) -> str:
        return self.content_service.build_report_title()

    def _build_summary_heading(self, request: AnalysisExportRequest) -> str:
        return self.content_service.build_summary_heading(request)

    def _build_representative_heading(self) -> str:
        return self.content_service.build_representative_heading()

    def _build_representative_sections(self, request: AnalysisExportRequest) -> list[tuple[str, list[str]]]:
        return self.content_service.build_representative_sections(request)

    def _build_ngram_representative_sections(self, request: AnalysisExportRequest) -> list[tuple[str, list[str]]]:
        return self.content_service.build_ngram_representative_sections(request)

    def _build_filename_stem(self, request: AnalysisExportRequest) -> str:
        return self.content_service.build_filename_stem(request)

    def _filters_text(self, request: AnalysisExportRequest) -> str:
        return self.content_service.filters_text(request)

    def _row_count_text(self, request: AnalysisExportRequest) -> str:
        return self.content_service.row_count_text(request)

    def _decode_chart(self, chart: AnalysisExportChartModel) -> DecodedChartImage:
        return self.chart_image_service.decode_chart(chart)

    def _fit_image_to_bounds(self, *, width: int, height: int, max_width: float, max_height: float) -> tuple[float, float]:
        return self.chart_image_service.fit_image_to_bounds(
            width=width,
            height=height,
            max_width=max_width,
            max_height=max_height,
        )

    def _display_column_label(self, value: str) -> str:
        return self.content_service.display_column_label(value)

    def _strip_extension(self, filename: str) -> str:
        return self.content_service.strip_extension(filename)

    def _slugify(self, value: str) -> str:
        return self.content_service.slugify(value)

    def _escape(self, value: str) -> str:
        return self.content_service.escape(value)

    def _truncate_text(self, value: str, *, limit: int) -> str:
        return self.content_service.truncate_text(value, limit=limit)

    def _sanitize_chart_caption(self, value: str | None) -> str:
        return self.content_service.sanitize_chart_caption(value)

    def _normalize_export_chart_image(self, chart: DecodedChartImage) -> DecodedChartImage:
        return self.chart_image_service.normalize_export_chart_image(chart)

    def _trim_pptx_chart_image(self, chart: DecodedChartImage) -> DecodedChartImage:
        return self.chart_image_service.trim_pptx_chart_image(chart)

    def _set_pptx_slide_background(self, slide) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = PptxRGBColor(*_PPTX_SLIDE_BACKGROUND_RGB)

    def _add_pptx_slide_title(self, slide, title: str) -> None:
        title_box = slide.shapes.add_textbox(
            PptxInches(_PPTX_CONTENT_LEFT),
            PptxInches(0.45),
            PptxInches(_PPTX_CONTENT_WIDTH),
            PptxInches(0.55),
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        paragraph = title_frame.paragraphs[0]
        paragraph.text = title
        paragraph.font.size = PptxPt(20)
        paragraph.font.bold = True
        paragraph.font.color.rgb = PptxRGBColor(*_PPTX_TEXT_RGB)
        paragraph.alignment = PP_ALIGN.LEFT
