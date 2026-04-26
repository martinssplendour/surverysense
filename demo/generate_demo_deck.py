from __future__ import annotations

import subprocess
from dataclasses import dataclass
from pathlib import Path, PurePosixPath

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT_PATH = ROOT / "demo" / "verbatim_demo_deck_2026-04-20.pptx"
BASELINE_REF = "91a7e58"
BASELINE_LABEL = "Friday April 17, 2026"
TODAY_LABEL = "Monday April 20, 2026"

WHITE = "FFFFFF"
BLACK = "111111"
GRAY = "666666"
LIGHT_GRAY = "F3F4F6"


@dataclass(frozen=True)
class Snapshot:
    label: str
    ref: str
    commit: str
    stats: dict[str, int]


def git(*args: str) -> str:
    return subprocess.check_output(["git", *args], cwd=ROOT, text=True).strip()


def tree_paths(ref: str) -> list[PurePosixPath]:
    listing = git("ls-tree", "-r", "--name-only", ref)
    return [PurePosixPath(line) for line in listing.splitlines() if line.strip()]


def collect_stats(ref: str) -> dict[str, int]:
    paths = tree_paths(ref)

    def count(predicate) -> int:
        return sum(1 for path in paths if predicate(path))

    return {
        "backend_py_modules": count(
            lambda p: p.suffix == ".py" and p.parts[:2] == ("backend", "app") and p.name != "__init__.py"
        ),
        "backend_api_modules": count(
            lambda p: p.suffix == ".py" and p.parts[:3] == ("backend", "app", "api") and p.name != "__init__.py"
        ),
        "backend_service_modules": count(
            lambda p: p.suffix == ".py" and p.parts[:3] == ("backend", "app", "services") and p.name != "__init__.py"
        ),
        "backend_service_packages": len(
            {
                p.parts[3]
                for p in paths
                if len(p.parts) > 4 and p.parts[:3] == ("backend", "app", "services")
            }
        ),
        "backend_test_modules": count(
            lambda p: p.suffix == ".py" and p.parts[:2] == ("backend", "tests") and p.name != "__init__.py"
        ),
        "frontend_js_modules": count(lambda p: p.suffix == ".js" and p.parts[0] == "frontend"),
        "frontend_results_modules": count(lambda p: p.suffix == ".js" and p.parts[:2] == ("frontend", "results")),
        "frontend_test_modules": count(
            lambda p: p.suffix in {".js", ".ts"} and p.parts[0] == "frontend" and ".test." in p.name
        ),
        "html_pages": count(lambda p: p.suffix == ".html" and p.parts[0] == "frontend"),
    }


def line_count(ref: str, path: str) -> int:
    return len(git("show", f"{ref}:{path}").splitlines())


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def set_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(WHITE)


def add_text(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    color: str = BLACK,
    font_name: str = "Aptos",
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    if paragraph.runs:
        run = paragraph.runs[0]
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)


def add_bullets(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    items: list[str],
    *,
    size: int = 16,
    color: str = BLACK,
) -> None:
    textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = f"- {item}"
        paragraph.alignment = PP_ALIGN.LEFT
        if paragraph.runs:
            run = paragraph.runs[0]
            run.font.name = "Aptos"
            run.font.size = Pt(size)
            run.font.color.rgb = rgb(color)
        paragraph.space_after = Pt(6)


def add_slide_title(slide, title: str, subtitle: str | None = None) -> None:
    add_text(slide, 0.7, 0.55, 12.0, 0.45, title, size=26, bold=True, font_name="Aptos Display")
    if subtitle:
        add_text(slide, 0.7, 1.0, 12.0, 0.28, subtitle, size=11, color=GRAY)


def add_section_heading(slide, left: float, top: float, text: str) -> None:
    add_text(slide, left, top, 5.5, 0.25, text, size=14, bold=True, color=BLACK)


def build_title_slide(prs: Presentation, friday: Snapshot, today: Snapshot) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(slide, "Verbatim App Demo Deck")
    add_text(
        slide,
        0.7,
        1.25,
        12.0,
        0.35,
        f"Baseline: {friday.label} ({friday.commit})   Current: {today.label} ({today.commit})",
        size=13,
        color=GRAY,
    )
    add_bullets(
        slide,
        0.9,
        2.0,
        8.0,
        2.4,
        [
            "Data ingestion and cleaning rationale",
            "Live demo flow",
            "Repo structure on Friday April 17, 2026",
            "Improvements through Monday April 20, 2026",
        ],
        size=18,
    )
    add_section_heading(slide, 0.7, 4.8, "Current snapshot")
    add_bullets(
        slide,
        0.9,
        5.1,
        8.0,
        1.4,
        [
            "1 Render web service",
            f'{today.stats["backend_py_modules"]} backend app modules',
            f'{today.stats["backend_service_modules"]} backend service modules',
            f'{today.stats["frontend_js_modules"]} frontend JS modules',
        ],
        size=16,
    )


def build_data_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(
        slide,
        "1. Data Ingestion and Cleaning",
        "Why this is engineered rather than treated as a single CSV import step.",
    )

    add_section_heading(slide, 0.7, 1.6, "Pipeline")
    add_bullets(
        slide,
        0.9,
        1.9,
        5.6,
        3.0,
        [
            "Read upload: detect encoding, parse CSV safely, and create preview versus architect samples.",
            "Diagnose layout: classify the file as wide or vertical using Gemini or rule-based heuristics.",
            "Transform deterministically: apply the manifest, clean headers, scrub null-like values, and enforce row limits.",
            "Build analysis dataset: keep metadata filterable and select true verbatim columns for NLP.",
            "Store result in memory for paging, filtering, reruns, representative examples, and export.",
        ],
        size=14,
    )

    add_section_heading(slide, 6.7, 1.6, "Why the engineering is justified")
    add_bullets(
        slide,
        6.9,
        1.9,
        5.6,
        3.0,
        [
            "Survey exports arrive in different shapes, not one standard table.",
            "Identifier columns, score columns, and fixed-response text must not be treated as open text.",
            "Multipart questions and duplicate answers need deterministic consolidation.",
            "The cleaned output must support interactive filtering and export, not only one-off analysis.",
        ],
        size=14,
    )

    add_section_heading(slide, 0.7, 5.35, "Covered in tests")
    add_bullets(
        slide,
        0.9,
        5.65,
        11.4,
        1.1,
        [
            "Vertical pivots, multipart word slots, UUID response ids, duplicate answer resolution, fixed-response rejection, and metadata detection.",
        ],
        size=14,
    )


def build_demo_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(
        slide,
        "2. Demo Flow",
        "A simple live path that shows both the product and the engineering story.",
    )

    add_section_heading(slide, 0.7, 1.6, "Recommended live path")
    add_bullets(
        slide,
        0.9,
        1.9,
        7.0,
        3.6,
        [
            "Upload a messy CSV and show row count, encoding, and preview rows.",
            "Open the manifest and call out wide versus vertical diagnosis plus metadata and verbatim detection.",
            "Show the transformed table to prove the cleaning path.",
            "Apply metadata filters to show filtered data and rerun behavior.",
            "Run BERTopic or K-means and open representative documents.",
            "Export PDF or Slides to show the workflow is end-to-end.",
        ],
        size=15,
    )

    add_section_heading(slide, 8.0, 1.6, "What to emphasise")
    add_bullets(
        slide,
        8.2,
        1.9,
        4.2,
        2.2,
        [
            "Same uploaded file flows through detection, transformation, filtering, analysis, and export.",
            "Show transformed data before model output.",
            "Use filters to show that the app is interactive, not static.",
        ],
        size=14,
    )

    add_section_heading(slide, 8.0, 4.45, "Best demo file")
    add_bullets(
        slide,
        8.2,
        4.75,
        4.2,
        1.2,
        [
            "Use a CSV with obvious metadata, at least one open-text column, and enough rows to make filtering visible.",
        ],
        size=14,
    )


def build_friday_slide(prs: Presentation, friday: Snapshot, lines: dict[str, tuple[int, int]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(
        slide,
        "3. Repo Structure on Friday",
        f"Snapshot from {friday.label}, commit {friday.commit}.",
    )

    add_section_heading(slide, 0.7, 1.6, "Friday structure")
    add_bullets(
        slide,
        0.9,
        1.9,
        5.7,
        2.5,
        [
            "1 runtime web service",
            f'{friday.stats["backend_py_modules"]} backend app modules',
            f'{friday.stats["backend_service_modules"]} backend service modules across {friday.stats["backend_service_packages"]} service packages',
            f'{friday.stats["backend_api_modules"]} API modules',
            f'{friday.stats["frontend_js_modules"]} frontend JS modules',
            f'{friday.stats["frontend_results_modules"]} results-page JS modules',
            f'{friday.stats["backend_test_modules"]} backend test modules',
        ],
        size=14,
    )

    add_section_heading(slide, 6.8, 1.6, "Architecture and state")
    add_bullets(
        slide,
        7.0,
        1.9,
        5.3,
        2.8,
        [
            "Static frontend served by FastAPI.",
            "routes_ingest.py handled upload, analysis, export, paging, and translation.",
            "Architect, cleaning, topic analysis, export, and result-store domains already existed.",
            "Uploaded results were kept in memory; there was no durable database.",
        ],
        size=14,
    )

    add_section_heading(slide, 0.7, 5.0, "Code shape on Friday")
    add_bullets(
        slide,
        0.9,
        5.3,
        11.3,
        1.2,
        [
            f'routes_ingest.py: {lines["routes_ingest.py"][0]} lines',
            f'workspace.js: {lines["workspace.js"][0]} lines',
            f'charts.js: {lines["charts.js"][0]} lines',
            f'analysis.js: {lines["analysis.js"][0]} lines',
        ],
        size=14,
    )


def build_improvements_slide(prs: Presentation, friday: Snapshot, today: Snapshot, lines: dict[str, tuple[int, int]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_slide_title(
        slide,
        "4. Improvements Between Friday and Today",
        f"Comparison: {BASELINE_LABEL} ({friday.commit}) versus {TODAY_LABEL} ({today.commit}).",
    )

    add_section_heading(slide, 0.7, 1.5, "Structure comparison")
    table = slide.shapes.add_table(8, 3, Inches(0.75), Inches(1.85), Inches(6.4), Inches(3.2)).table
    headers = ["Metric", "Friday", "Today"]
    for column_index, value in enumerate(headers):
        cell = table.cell(0, column_index)
        cell.text = value
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb(LIGHT_GRAY)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.color.rgb = rgb(BLACK)

    rows = [
        ("Runtime web services", "1", "1"),
        ("Backend app modules", str(friday.stats["backend_py_modules"]), str(today.stats["backend_py_modules"])),
        (
            "Backend service modules",
            f'{friday.stats["backend_service_modules"]} in {friday.stats["backend_service_packages"]} packages',
            f'{today.stats["backend_service_modules"]} in {today.stats["backend_service_packages"]} packages',
        ),
        ("API modules", str(friday.stats["backend_api_modules"]), str(today.stats["backend_api_modules"])),
        ("Frontend JS modules", str(friday.stats["frontend_js_modules"]), str(today.stats["frontend_js_modules"])),
        ("Results JS modules", str(friday.stats["frontend_results_modules"]), str(today.stats["frontend_results_modules"])),
        ("Backend test modules", str(friday.stats["backend_test_modules"]), str(today.stats["backend_test_modules"])),
    ]
    for row_index, row in enumerate(rows, start=1):
        for column_index, value in enumerate(row):
            cell = table.cell(row_index, column_index)
            cell.text = value
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)
                    run.font.color.rgb = rgb(BLACK)

    add_section_heading(slide, 7.5, 1.5, "What improved")
    add_bullets(
        slide,
        7.7,
        1.85,
        4.6,
        2.4,
        [
            f'routes_ingest.py: {lines["routes_ingest.py"][0]} -> {lines["routes_ingest.py"][1]} lines',
            f'report_export_service.py: {lines["report_export_service.py"][0]} -> {lines["report_export_service.py"][1]} lines',
            f'topic_analysis_service.py: {lines["topic_analysis_service.py"][0]} -> {lines["topic_analysis_service.py"][1]} lines',
            f'main.py: {lines["main.py"][0]} -> {lines["main.py"][1]} lines',
            f'workspace.js: {lines["workspace.js"][0]} -> {lines["workspace.js"][1]} lines',
            f'charts.js: {lines["charts.js"][0]} -> {lines["charts.js"][1]} lines',
            f'analysis.js: {lines["analysis.js"][0]} -> {lines["analysis.js"][1]} lines',
        ],
        size=12,
    )

    add_section_heading(slide, 0.7, 5.4, "Why this matters")
    add_bullets(
        slide,
        0.9,
        5.7,
        11.4,
        1.0,
        [
            "The module count increased because responsibilities were split into smaller units; the deployed system did not add more runtime services.",
            "Application wiring is now explicit in application_setup.py.",
            "Ingest routes are split into upload, analysis, result, and translation modules with shared context.",
            "Frontend state is explicit in state.js, with workspace, persistence, chart, and analysis behavior split into focused files.",
        ],
        size=13,
    )


def main() -> None:
    today_ref = "HEAD"
    friday = Snapshot(
        label=BASELINE_LABEL,
        ref=BASELINE_REF,
        commit=git("rev-parse", "--short", BASELINE_REF),
        stats=collect_stats(BASELINE_REF),
    )
    today = Snapshot(
        label=TODAY_LABEL,
        ref=today_ref,
        commit=git("rev-parse", "--short", today_ref),
        stats=collect_stats(today_ref),
    )
    lines = {
        "routes_ingest.py": (
            line_count(BASELINE_REF, "backend/app/api/routes_ingest.py"),
            line_count(today_ref, "backend/app/api/routes_ingest.py"),
        ),
        "report_export_service.py": (
            line_count(BASELINE_REF, "backend/app/services/report_export_service/report_export_service.py"),
            line_count(today_ref, "backend/app/services/report_export_service/report_export_service.py"),
        ),
        "topic_analysis_service.py": (
            line_count(BASELINE_REF, "backend/app/services/topic_analysis_services/topic_analysis_service.py"),
            line_count(today_ref, "backend/app/services/topic_analysis_services/topic_analysis_service.py"),
        ),
        "workspace.js": (
            line_count(BASELINE_REF, "frontend/results/workspace.js"),
            line_count(today_ref, "frontend/results/workspace.js"),
        ),
        "charts.js": (
            line_count(BASELINE_REF, "frontend/results/charts.js"),
            line_count(today_ref, "frontend/results/charts.js"),
        ),
        "analysis.js": (
            line_count(BASELINE_REF, "frontend/results/analysis.js"),
            line_count(today_ref, "frontend/results/analysis.js"),
        ),
        "main.py": (
            line_count(BASELINE_REF, "backend/app/main.py"),
            line_count(today_ref, "backend/app/main.py"),
        ),
    }

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_title_slide(prs, friday, today)
    build_data_slide(prs)
    build_demo_slide(prs)
    build_friday_slide(prs, friday, lines)
    build_improvements_slide(prs, friday, today, lines)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(OUTPUT_PATH)


if __name__ == "__main__":
    main()
